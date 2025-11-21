# TITLE: PPTX to PDF
# AUTHOR: @sebastian

from pptx import Presentation
from reportlab.lib.pagesizes import letter, landscape
from reportlab.pdfgen import canvas

title = "PPTX to PDF"
input_format = "pptx"
output_format = "pdf"


def convert(file_path):
    new_path = file_path.replace(".pptx", "_converted.pdf")
    
    prs = Presentation(file_path)
    c = canvas.Canvas(new_path, pagesize=landscape(letter))
    width, height = landscape(letter)
    
    for slide in prs.slides:
        y_position = height - 100
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                c.drawString(100, y_position, shape.text[:100])
                y_position -= 20
        
        c.showPage()
    
    c.save()
    return new_path